"""Builds the 10-slide executive PowerPoint deck at
outputs/meridian_home_deck.pptx using python-pptx. Dark navy theme with
amber accents, populated with real numbers from the scoring and financial
models.
"""

import json
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PROJECT_ROOT = Path(__file__).parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

RAW_DIR = PROJECT_ROOT / "data_collection" / "raw"
CHARTS_DIR = PROJECT_ROOT / "outputs" / "charts"
OUTPUTS_DIR = PROJECT_ROOT / "outputs"

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xE8, 0xEA, 0xF6)
AMBER = RGBColor(0xF9, 0xA8, 0x25)
GREEN = RGBColor(0x0C, 0xA3, 0x0C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    return slide


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _add_title(slide, text):
    return _add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8), text, size=30, bold=True, color=WHITE)


def _add_divider(slide, top=Inches(1.15)):
    line = slide.shapes.add_shape(1, Inches(0.6), top, Inches(2.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = AMBER
    line.line.fill.background()


def _add_bullets(slide, left, top, width, height, bullets, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = Pt(size)
        p.font.color.rgb = LIGHT_GREY
        p.font.name = "Calibri"
        p.space_after = Pt(14)
    return box


def _add_picture_fit(slide, image_path, left, top, max_width, max_height):
    from PIL import Image

    with Image.open(image_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h
    width = max_width
    height = Emu(int(width / aspect))
    if height > max_height:
        height = max_height
        width = Emu(int(height * aspect))
    centered_left = Emu(int(left + (max_width - width) / 2))
    slide.shapes.add_picture(str(image_path), centered_left, top, width=width, height=height)


def build_deck():
    scores = pd.read_csv(RAW_DIR / "city_scores.csv").sort_values("rank")
    financial_results = json.loads((RAW_DIR / "financial_results.json").read_text())

    top2 = scores[scores["recommended"]]
    city1 = top2.iloc[0]  # Indianapolis
    city2 = top2.iloc[1]  # Charlotte
    fin1 = financial_results[city1["metro_name"]]
    fin2 = financial_results[city2["metro_name"]]

    combined_npv_risk = fin1["npv_risk_adjusted"] + fin2["npv_risk_adjusted"]
    combined_y3_ebitda = fin1["key_metrics"]["y3_ebitda"] + fin2["key_metrics"]["y3_ebitda"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- Slide 1: Title ---
    s = _blank_slide(prs)
    _add_textbox(s, Inches(1), Inches(2.4), Inches(11.3), Inches(1.2), "Meridian Home", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(1), Inches(3.4), Inches(11.3), Inches(0.7), "Market Expansion Strategy", size=24, color=AMBER, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6), "Which Two Markets Should We Enter First?", size=16, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    line = s.shapes.add_shape(1, Inches(5.667), Inches(3.25), Inches(2), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = AMBER; line.line.fill.background()
    _add_textbox(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.4), "Prepared by Balraj Kooner | Masters in Business Analytics", size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    _add_textbox(s, Inches(1), Inches(6.95), Inches(11.3), Inches(0.4), "July 2026", size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # --- Slide 2: Situation ---
    s = _blank_slide(prs)
    _add_title(s, "Situation")
    _add_divider(s)
    _add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(5), [
        "Meridian Home operates 47 stores concentrated in the Northeast and West Coast, with ~$890M in annual revenue and strong brand recognition in the premium home goods segment",
        "The board has approved expansion into 2 new markets, targeting metros with strong homeowner demographics and limited premium competition",
        "Five candidate metros were identified across three growth regions: Sun Belt (Nashville, Austin, Charlotte), Mountain West (Denver), Midwest (Indianapolis)",
    ], size=18)

    # --- Slide 3: Complication ---
    s = _blank_slide(prs)
    _add_title(s, "Complication")
    _add_divider(s)
    _add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(5), [
        "Not all five markets offer equivalent returns — market attractiveness, competitive density, and unit economics vary significantly across the candidates",
        "Entering the wrong market carries substantial risk: $450K pre-opening investment per store, multi-year lease commitments, and reputational risk of a poorly-performing flagship in a new region",
        "A rigorous, data-driven framework is needed to identify the markets with the highest probability of success across 7 weighted dimensions — and to be honest about where the unit economics are still marginal",
    ], size=18)

    # --- Slide 4: Recommendation ---
    s = _blank_slide(prs)
    _add_title(s, "Our Recommendation")
    _add_divider(s)
    box = s.shapes.add_shape(1, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = AMBER; box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Enter Charlotte first (Q1 2026); open a single Indianapolis pilot store in parallel"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    _add_bullets(s, Inches(0.7), Inches(2.8), Inches(11.8), Inches(4.3), [
        f"Indianapolis ranks #1 on market attractiveness (composite {city1['composite_score']:.2f}/5.0) — cheapest lease rate (${city1['retail_lease_rate_sqft_annual']:.2f}/sqft) and lowest competitive density of any candidate — but slower population growth ({city1['population_growth_rate_pct']:.1f}%/5yr) means its financial case needs more runway to prove out than our 3-year model window",
        f"Charlotte ranks #2 (composite {city2['composite_score']:.2f}/5.0), sits at Meridian's own East Coast distribution hub (zero logistics cost), and is the only recommended market that clears payback within the model horizon (~{fin2['payback_months']:.0f} months)",
        f"Combined risk-adjusted 3-year NPV is currently negative (${combined_npv_risk:,.0f}) at the standard 4,500 sqft / 10-FTE format — this is a real constraint, not a rounding error, and the phased approach below is how we manage that risk",
        "Recommended sequencing: open Charlotte's 2-store cluster now; open 1 Indianapolis pilot store (not 2) to validate real-world performance before committing the second store there",
    ], size=17)

    # --- Slide 5: Methodology ---
    s = _blank_slide(prs)
    _add_title(s, "Analytical Framework")
    _add_divider(s)
    _add_textbox(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "Issue Tree", size=16, bold=True, color=AMBER)
    _add_bullets(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(4.8), [
        "Is the market attractive? (size, income, housing, e-commerce penetration)",
        "Can Meridian compete effectively? (saturation, cannibalization, logistics)",
        "What is the financial return? (revenue, cost, breakeven, NPV)",
    ], size=15)

    _add_textbox(s, Inches(6.7), Inches(1.5), Inches(5.9), Inches(0.4), "7 Evaluation Dimensions", size=16, bold=True, color=AMBER)
    dims = [
        ("Population & Growth", "15%"), ("Income & Furnishing Spend", "20%"),
        ("Homeownership & Housing", "20%"), ("Competitive Saturation", "20%"),
        ("Retail Lease Cost", "10%"), ("Labor Cost", "10%"), ("Logistics Proximity", "5%"),
    ]
    top = Inches(2.0)
    for name, weight in dims:
        _add_textbox(s, Inches(6.7), top, Inches(4.6), Inches(0.4), name, size=14, color=LIGHT_GREY)
        _add_textbox(s, Inches(11.2), top, Inches(1.2), Inches(0.4), weight, size=14, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
        top += Emu(int(Inches(0.5)))

    # --- Slide 6: Market attractiveness comparison ---
    s = _blank_slide(prs)
    _add_title(s, "Market Attractiveness Scorecard")
    _add_divider(s)
    _add_picture_fit(s, CHARTS_DIR / "composite_scores.png", Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.3))
    _add_textbox(s, Inches(0.7), Inches(6.95), Inches(11.9), Inches(0.4),
                 "Source: US Census ACS, BLS OES/CEX, Zillow Research, NAHB, CommercialCafe — analysis by Balraj Kooner",
                 size=10, color=LIGHT_GREY)

    # --- Slide 7: Competitive landscape ---
    s = _blank_slide(prs)
    _add_title(s, "Competitive Saturation Analysis")
    _add_divider(s)
    _add_picture_fit(s, CHARTS_DIR / "competitive_heatmap.png", Inches(0.7), Inches(1.5), Inches(11.9), Inches(4.3))
    least_saturated = scores.sort_values("competitors_per_100k").iloc[0]
    _add_textbox(
        s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(1.2),
        f"Indianapolis and Austin show the lowest direct competitor density (~0.21-0.24 per 100k population) — "
        f"Charlotte offers a reasonable balance of growth and moderate competition (0.29 per 100k)",
        size=15, color=AMBER,
    )

    # --- Slide 8: Financial model ---
    s = _blank_slide(prs)
    _add_title(s, "3-Year Financial Projection — Recommended Markets")
    _add_divider(s)
    _add_picture_fit(s, CHARTS_DIR / "financial_waterfall.png", Inches(0.7), Inches(1.4), Inches(11.9), Inches(4.4))

    metrics = [
        ("Combined Year 3 EBITDA", f"${combined_y3_ebitda:,.0f}"),
        ("Charlotte Payback Period", f"{fin2['payback_months']:.0f} months"),
        ("Combined Risk-Adj. NPV", f"${combined_npv_risk:,.0f}"),
    ]
    col_width = Inches(3.9)
    for i, (label, value) in enumerate(metrics):
        left = Inches(0.7) + Emu(int(col_width) * i)
        _add_textbox(s, left, Inches(6.0), col_width, Inches(0.4), label, size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
        _add_textbox(s, left, Inches(6.4), col_width, Inches(0.5), value, size=20, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    # --- Slide 9: Risks and mitigation ---
    s = _blank_slide(prs)
    _add_title(s, "Key Risks and Mitigation")
    _add_divider(s)
    risks = [
        ("Combined NPV is negative at the modeled 4,500 sqft / 10-FTE format",
         "High",
         "Phase Indianapolis as a single pilot store; revisit store format (sqft, staffing) before the second store"),
        ("Competitive response from West Elm or Pottery Barn",
         "Medium",
         "Pre-commit to Charlotte's 2-store cluster to establish critical mass before competitors react"),
        ("Ramp-up slower than modeled (consumer awareness lag)",
         "Medium",
         "Year 1 budget includes incremental local marketing spend beyond the 3%-of-revenue baseline"),
    ]
    top = Inches(1.6)
    col_widths = [Inches(5.3), Inches(1.4), Inches(5.2)]
    headers = ["Risk", "Probability", "Mitigation"]
    left = Inches(0.7)
    for header, w in zip(headers, col_widths):
        _add_textbox(s, left, top, w, Inches(0.4), header, size=14, bold=True, color=AMBER)
        left = Emu(int(left) + int(w))
    top = Inches(2.1)
    for risk, prob, mitigation in risks:
        left = Inches(0.7)
        for text, w in zip([risk, prob, mitigation], col_widths):
            _add_textbox(s, left, top, w, Inches(1.3), text, size=13, color=LIGHT_GREY)
            left = Emu(int(left) + int(w))
        top = Emu(int(top) + int(Inches(1.5)))

    # --- Slide 10: Implementation roadmap ---
    s = _blank_slide(prs)
    _add_title(s, "18-Month Implementation Roadmap")
    _add_divider(s)
    roadmap = [
        "Month 1-2: Market validation (real estate broker engagement, trade area confirmation) — both markets",
        "Month 2-4: Charlotte site selection and LOI negotiation (2-store cluster)",
        "Month 4-6: Charlotte lease execution and buildout",
        "Month 6: Charlotte stores open — Q1 2026",
        "Month 6-8: Indianapolis single-store site selection and LOI (pilot format)",
        "Month 8-10: Indianapolis pilot buildout",
        "Month 10: Indianapolis pilot store opens",
        "Month 10-18: Performance monitoring against model; go/no-go decision on Indianapolis store #2",
    ]
    _add_bullets(s, Inches(0.7), Inches(1.6), Inches(11.8), Inches(5.3), roadmap, size=16)

    out_path = OUTPUTS_DIR / "meridian_home_deck.pptx"
    prs.save(out_path)
    print(f"Saved {out_path} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build_deck()
